import html
import os
import subprocess
import tempfile
from io import BytesIO
from tempfile import SpooledTemporaryFile
from typing import Dict, List, Tuple

import fitz
import pandas as pd
from azure.ai.formrecognizer import DocumentAnalysisClient
from docx import Document
from fastapi import UploadFile
from pptx import Presentation
from pypdf import PdfReader

from app.core.config import settings


def read_docx(docx_file: UploadFile) -> str:
    """Reads a Docx file a returns a str of convertable text
    does not require LibreOffice. Cannot convert Image to text

    Args:
        docx_file (UploadFile): The DOCX file to convert.
    Returns:
        (str): a string of all convertable text in the document
    """
    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as temp:
        docx_file.file.seek(0)
        data = docx_file.file.read()
        temp.write(data)
        temp_filename = temp.name
        # Use python-docx to read the temporary file
        doc = Document(temp_filename)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])

    return text


def read_pptx(pptx_file: UploadFile) -> str:
    """Reads a PPTX file a returns a str of convertable text
    does not require LibreOffice. Cannot convert Image to text

    Args:
        pptx_file (UploadFile): The DOCX file to convert.
    Returns:
        (str): a string of all convertable text in the document
    """
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp:
        pptx_file.file.seek(0)
        contents = pptx_file.file.read()
        temp.write(contents)
        temp_filename = temp.name
        # Use python-pptx to read the temporary file
        presentation = Presentation(temp_filename)

        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "

    return text


def doc_to_pdf(file: UploadFile) -> UploadFile:
    """
    Converts an uploaded document file to PDF using LibreOffice and
    returns the PDF as an UploadFile. Requires LibreOffice running in
    background as a headless daemon.

    Args:
        file (UploadFile): The file to convert.

    Returns:
        UploadFile: An UploadFile object of the converted PDF.

    Raises:
        Exception: If LibreOffice is not installed.
    """
    try:
        # Check if LibreOffice is installed
        subprocess.run(["libreoffice", "--version"], check=True)
    except subprocess.CalledProcessError:
        raise Exception(
            user_message=f"Failed to read file {file.filename}, "
            "please try another file.",
            internal_logging_message="Please install LibreOffice to doc "
            "slice PPTX & DOCX.",
        )

    # writing the pptx / docx file to a temp file object
    with tempfile.NamedTemporaryFile() as temp_file:
        temp_file_name = temp_file.name
        with open(temp_file_name, "wb") as buffer:
            file.file.seek(0)
            buffer.write(file.file.read())

        pdf_file_name = os.path.splitext(temp_file_name)[0] + ".pdf"

        # Convert pptx / docx to pdf using libreoffice
        command = [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            temp_file_name,
            "--outdir",
            os.path.dirname(temp_file_name),
        ]
        subprocess.run(command, check=True)

    # Create an SpooledTemporaryFile (which can be used by UploadFile)
    pdf_file = SpooledTemporaryFile()
    with open(pdf_file_name, "rb") as f:
        f.seek(0)
        pdf_file.write(f.read())
    pdf_file.seek(0)

    # Remove the temporary PDF file
    os.remove(pdf_file_name)

    # Create and return an UploadFile object
    return UploadFile(BytesIO(pdf_file.read()), filename=file.filename)


def pdf_to_page_map_azure(file: UploadFile) -> List[Tuple[int, int, str]]:
    """
    takes a pdf document and turn it into a page map,
    using Azure Form Recognizer and DocumentAnalysisClient
    by mapping each page into a list, creating a page map

    Args:
        file (UploadFile): The file to be unpacked.

    Returns:
        page_map (list[Tuple[int, int, str]]): a list of pages to be read
            eg: [(int, int, str)]
            eg: [({page num base 0},
            {character position where it starts}, {text content})]

    """

    def _table_to_html(table) -> str:
        """
        function within get_document_text,
        turns a table into str, html that's readable

        args:
            table

        return:
            (str): table in html format
        """
        table_html = "<table>"
        rows = [
            sorted(
                [cell for cell in table.cells if cell.row_index == i],
                key=lambda cell: cell.column_index,
            )
            for i in range(table.row_count)
        ]
        for row_cells in rows:
            table_html += "<tr>"
            for cell in row_cells:
                tag = (
                    "th"
                    if (cell.kind == "columnHeader" or cell.kind == "rowHeader")
                    else "td"
                )
                cell_spans = ""
                if cell.column_span > 1:
                    cell_spans += f" colSpan={cell.column_span}"
                if cell.row_span > 1:
                    cell_spans += f" rowSpan={cell.row_span}"
                table_html += f"<{tag}{cell_spans}>{html.escape(cell.content)}</{tag}>"
            table_html += "</tr>"
        table_html += "</table>"
        return table_html

    offset = 0
    page_map: List[Tuple[int, int, str]] = []

    if settings.FORMRECOGNIZER_CREDS is None:
        raise Exception(
            user_message="Please contact administrator to enable additional features",
            internal_logging_message="Form recognizer credentials are required.",
        )
    form_recognizer_client = DocumentAnalysisClient(
        endpoint=settings.FORM_RECOGNIZER_SERVICE,
        credential=settings.FORMRECOGNIZER_CREDS,
        headers={"x-ms-useragent": "azure-search-chat-demo/1.0.0"},
    )

    # Read file data into bytes, file.file.seek(0) is important!
    file.file.seek(0)
    file_data = file.file.read()

    # Create a BytesIO object from the file data
    file_stream = BytesIO(file_data)

    poller = form_recognizer_client.begin_analyze_document(
        "prebuilt-layout", document=file_stream
    )
    form_recognizer_results = poller.result()
    # await file.close()

    for page_num, page in enumerate(form_recognizer_results.pages):
        tables_on_page = []
        if form_recognizer_results.tables is not None:
            for table in form_recognizer_results.tables:
                if table.bounding_regions is None:
                    continue
                if table.bounding_regions[0].page_number == page_num + 1:
                    tables_on_page.append(table)

        # mark all positions of the table spans in the page
        page_offset = page.spans[0].offset
        page_length = page.spans[0].length
        table_chars = [-1] * page_length
        for table_id, table in enumerate(tables_on_page):
            for span in table.spans:
                # replace all table spans with "table_id" in table_chars array
                for i in range(span.length):
                    idx = span.offset - page_offset + i
                    if idx >= 0 and idx < page_length:
                        table_chars[idx] = table_id

        # build page text by replacing charcters in table spans with table html
        page_text = ""
        added_tables = set()
        for idx, table_id in enumerate(table_chars):
            if table_id == -1:
                page_text += form_recognizer_results.content[page_offset + idx]
            elif table_id not in added_tables:
                page_text += _table_to_html(tables_on_page[table_id])
                added_tables.add(table_id)

        page_text += " "
        page_map.append((page_num, offset, page_text))
        offset += len(page_text)

    return page_map


def pdf_to_page_map_pypdf(file: UploadFile) -> List[Tuple[int, int, str]]:
    """
    takes a pdf document and turn it into a page map,
    using pypdf to perform non-OCR extraction
    by mapping each page into a list, creating a page map

    Args:
        file (UploadFile): The file to be unpacked.

    Returns:
        page_map (list[Tuple[int, int, str]]): a list of pages to be read
            eg: [(int, int, str)]
            eg: [({page num base 0},
            {character position where it starts}, {text content})]

    """
    file.file.seek(0)
    file_data = file.file.read()

    # Create a BytesIO object from the file data
    file_stream = BytesIO(file_data)
    reader = PdfReader(file_stream)

    pages_text = []
    for page_num, page in enumerate(reader.pages):
        text = page.extract_text() if page.extract_text() is not None else ""
        pages_text.append((page_num, 0, text))

    return pages_text


def pdf_to_page_map_pymupdf(file: UploadFile) -> List[Tuple[int, int, str]]:
    """
    takes a pdf document and turn it into a page map,
    using PyMuPDF (aka Pyfitz) to perform non-OCR extraction
    by mapping each page into a list, creating a page map

    Args:
        file (UploadFile): The file to be unpacked.

    Returns:
        page_map (list[Tuple[int, int, str]]): a list of pages to be read
            eg: [(int, int, str)]
            eg: [({page num base 0},
            {character position where it starts}, {text content})]

    """
    file.file.seek(0)
    file_data = file.file.read()

    # Create a BytesIO object from the file data
    file_stream = BytesIO(file_data)

    pdf_document = fitz.open(stream=file_stream, filetype="pdf")

    pages_text = []
    for page_num in range(len(pdf_document)):
        page = pdf_document[page_num]
        text = page.get_text("text")
        pages_text.append((page_num, 0, text))

    pdf_document.close()

    return pages_text


def extract_csv(file: UploadFile, csv_header) -> List[Dict]:
    """
    convert the csv into a list of dict,
    while each dict representing a column, always return a dictionary of
    strings.

    args:
        file (UploadFile): the file to be extracted
        csv_header (bool): whether the csv contains a header

    return:
        (List[Dict]):
            a dict representation of a csv, if csv_header is True:
            [
                {"i_title": "Issue1", "s_title": "Solution1", "e_title": "Extra1"},
                {"i_title": "Issue2", "s_title": "Solution2", "e_title": "Extra2"},
                {"i_title": "Issue3", "s_title": "Solution3", "e_title": "Extra3"},
            ]
            if csv_header is False (kn are column names generated):
            [
                {"k0": "i_title", "k1": "s_title", "k2": "e_title"},
                {"k0": "Issue1", "k1": "Solution1", "k2": "Extra1"},
                {"k0": "Issue2", "k1": "Solution2", "k2": "Extra2"},
                {"k0": "Issue3", "k1": "Solution3", "k2": "Extra3"},
            ]
    """
    csvdf = pd.read_csv(file.file, encoding="utf-8", header=0 if csv_header else None)

    record = csvdf.to_dict(orient="records")

    # convert all int keys to string - 0 -> "k0" for mapping later
    return (
        record if csv_header else [{f"k{k}": v for k, v in d.items()} for d in record]
    )
